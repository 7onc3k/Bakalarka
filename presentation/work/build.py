"""Build obhajoba .pptx draft v1 from FIS template.

Strategy:
- Load template (8 slides). FIS layouts: 'Úvodní snímek' + 'Nadpis a obsah'.
- Insert a NEW slide between #5 and #6 (Výsledky II: ablace) via reorder of sldIdLst.
- Replace title + content placeholders on every slide.
- Add custom shapes/tables on slides 5, 6, 9.
- Speaker notes on every slide = mluvený text.

Output: obhajoba-draft-v1.pptx
"""
from copy import deepcopy
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Cm, Emu

HERE = Path(__file__).parent
TEMPLATE = HERE.parent / "templates" / "Sablona_informaticke_PROGRAMY.pptx"
OUT = HERE / "obhajoba-draft-v1.pptx"

# FIS green from template (#009881)
FIS_GREEN = RGBColor(0x00, 0x98, 0x81)
DARK_GRAY = RGBColor(0x4A, 0x4A, 0x49)
LIGHT_GRAY = RGBColor(0xE8, 0xE8, 0xE8)
ACCENT_RED = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_ORANGE = RGBColor(0xE6, 0x8A, 0x00)

# Title placeholder occupies y=1.20–1.95.
# Custom content must start at y >= 2.05 to avoid overlap.
CONTENT_TOP = Inches(2.05)


def remove_empty_placeholder(slide, idx):
    """Remove placeholder by idx if present (to make room for custom layout)."""
    for shp in list(slide.shapes):
        if shp.is_placeholder and shp.placeholder_format.idx == idx:
            sp = shp._element
            sp.getparent().remove(sp)
            return

# =============================================================================
# HELPERS
# =============================================================================

def replace_text_keep_format(shape, new_paragraphs):
    """Replace text in shape, keeping formatting of first run as template.

    new_paragraphs: list of either str OR list of (text, bold) tuples
    """
    tf = shape.text_frame
    first_p = tf.paragraphs[0]
    if first_p.runs:
        template_rPr = deepcopy(first_p.runs[0]._r.get_or_add_rPr())
    else:
        template_rPr = None

    # Clear all paragraphs after first
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    # Clear all runs in first paragraph
    for r in list(first_p.runs):
        r._r.getparent().remove(r._r)
    # Also remove any line breaks
    for br in first_p._p.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}br"
    ):
        first_p._p.remove(br)

    def add_run(p, text, bold=None):
        run = p.add_run()
        run.text = text
        if template_rPr is not None:
            new_rPr = deepcopy(template_rPr)
            run._r.insert(0, new_rPr)
            # Remove the auto-created rPr to keep our copy
            rPrs = run._r.findall(
                "{http://schemas.openxmlformats.org/drawingml/2006/main}rPr"
            )
            for rPr in rPrs[1:]:
                run._r.remove(rPr)
        if bold is not None:
            run.font.bold = bold

    for i, para in enumerate(new_paragraphs):
        if i == 0:
            p = first_p
        else:
            p = tf.add_paragraph()
        if isinstance(para, str):
            add_run(p, para)
        else:
            for chunk in para:
                if isinstance(chunk, str):
                    add_run(p, chunk)
                else:
                    text, bold = chunk
                    add_run(p, text, bold=bold)


def set_notes(slide, text):
    """Set speaker notes."""
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.clear()
    # Split into paragraphs by \n\n
    paras = text.strip().split("\n\n")
    for i, para in enumerate(paras):
        p = notes_tf.paragraphs[0] if i == 0 else notes_tf.add_paragraph()
        # Clear any existing runs
        for r in list(p.runs):
            r._r.getparent().remove(r._r)
        run = p.add_run()
        run.text = para
        run.font.size = Pt(11)


def add_textbox(slide, left, top, width, height, text_paragraphs,
                font_size=14, bold=False, color=DARK_GRAY, align=PP_ALIGN.LEFT,
                fill=None, border=None, italic=False):
    """Add a simple text box."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    if border is not None:
        tb.line.color.rgb = border
        tb.line.width = Pt(1)
    else:
        tb.line.fill.background()
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    for i, para in enumerate(text_paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        # Normalize paragraph → list of (text, bold, color, italic) 4-tuples.
        # Chunk forms allowed:
        #   - str
        #   - (text, bold)
        #   - (text, bold, color)
        #   - (text, bold, color, italic)
        if isinstance(para, str):
            chunks = [(para, bold, color, italic)]
        elif isinstance(para, tuple) and len(para) == 2 and isinstance(para[1], bool):
            chunks = [(para[0], para[1], color, italic)]
        else:
            chunks = []
            for chunk in para:
                if isinstance(chunk, str):
                    chunks.append((chunk, bold, color, italic))
                elif isinstance(chunk, tuple):
                    if len(chunk) == 4:
                        chunks.append(chunk)
                    elif len(chunk) == 3:
                        chunks.append((chunk[0], chunk[1], chunk[2], italic))
                    elif len(chunk) == 2:
                        chunks.append((chunk[0], chunk[1], color, italic))
        for chunk_text, chunk_bold, chunk_color, chunk_italic in chunks:
            run = p.add_run()
            run.text = chunk_text
            run.font.size = Pt(font_size)
            run.font.bold = chunk_bold
            run.font.italic = chunk_italic
            run.font.color.rgb = chunk_color
    return tb


def add_mini_table(slide, left, top, width, height, rows):
    """rows = list of lists, first row is header. Used by slide 10 backup."""
    cols = len(rows[0])
    tbl = slide.shapes.add_table(len(rows), cols, left, top, width, height).table
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(11)
            run.font.color.rgb = DARK_GRAY
            if i == 0:
                run.font.bold = True
                run.font.color.rgb = FIS_GREEN
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
    return tbl


def add_arrow_shape(slide, text, left, top, width, height, fill_color, font_size=14):
    """Add a chevron/arrow shape with text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return shape


def insert_new_slide_after(prs, after_index, layout):
    """Insert new slide at position after_index+1 by reordering sldIdLst."""
    new_slide = prs.slides.add_slide(layout)
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides = list(xml_slides)
    # Newly added slide is at the end
    new_sld_id = slides[-1]
    xml_slides.remove(new_sld_id)
    # Insert after `after_index`
    xml_slides.insert(after_index + 1, new_sld_id)
    return new_slide


def remove_slide(prs, slide_index):
    """Remove slide at given index. python-pptx doesn't expose this natively."""
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides_list = list(xml_slides)
    sld_id_elem = slides_list[slide_index]
    rId = sld_id_elem.attrib[
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    ]
    prs.part.drop_rel(rId)
    xml_slides.remove(sld_id_elem)


def reorder_slides_to(prs, target_order):
    """Reorder slides by current 0-based indices (target_order is list of current indices)."""
    xml_slides = prs.slides._sldIdLst  # noqa: SLF001
    slides = list(xml_slides)
    for s in slides:
        xml_slides.remove(s)
    for idx in target_order:
        xml_slides.append(slides[idx])


# =============================================================================
# BUILD
# =============================================================================

prs = Presentation(str(TEMPLATE))
# Template has 3 slide masters; the layouts we need are on master 2
all_layouts = []
for master in prs.slide_masters:
    for L in master.slide_layouts:
        all_layouts.append(L)
title_layout = next(L for L in all_layouts if L.name == "Úvodní snímek")
content_layout = next(L for L in all_layouts if L.name == "Nadpis a obsah")

# -----------------------------------------------------------------------------
# SLIDE 1 — Titul
# -----------------------------------------------------------------------------
s = prs.slides[0]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    txt = shp.text_frame.text
    if txt == "Název práce":
        replace_text_keep_format(shp, [
            "Návrh a testování běhových prostředí pro autonomní coding agenty v softwarovém vývoji"
        ])
    elif txt == "Vedoucí práce:":
        replace_text_keep_format(shp, ["Vedoucí práce: Ing. Jiří Korčák"])
    elif txt == "Jméno a příjmení autora":
        replace_text_keep_format(shp, ["Thanh An Nguyen"])
    elif txt == "Studijní program:":
        replace_text_keep_format(shp, ["Studijní program: Aplikovaná informatika"])
    elif "Specializace" in txt:
        # Specializace odstraněna — program AIN ji nemá; vyprázdnit text
        replace_text_keep_format(shp, [""])
    elif "Datum obhajoby" in txt:
        replace_text_keep_format(shp, ["Datum obhajoby: [TODO červen 2026]"])

set_notes(s, """[HOOK ~30 s]

Dobrý den. Když dnes dáte LLM agentovi spec a necháte ho napsat kód, ve většině případů ho napíše. Otázka, kterou si tato práce klade, je: proč zrovna takhle? A jestli to, jak ho píše, lze nějak řídit.

[Plynule navázat na slide 2.]""")

# -----------------------------------------------------------------------------
# SLIDE 2 — Téma a struktura prezentace (FIS template slot)
# -----------------------------------------------------------------------------
# Title placeholder rozměry: x=0.688 → x=9.195 (width 8.508). Content sloupce
# musí být zarovnané přesně na tyto okraje, jinak vizuálně přečnívají.
# Layout: motivace vlevo (5.0"), struktura vpravo (3.2"), gap 0.3".
# Text v bodech (kotvy), ne věty — řečník rozebere v mluveném textu.
s = prs.slides[1]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, ["Téma a struktura prezentace"])
    elif shp.placeholder_format and shp.placeholder_format.idx == 1:
        replace_text_keep_format(shp, [""])

remove_empty_placeholder(s, 1)

# Alignment k title placeholderu
LEFT_X, LEFT_W = Inches(0.688), Inches(5.0)
RIGHT_X, RIGHT_W = Inches(5.988), Inches(3.207)

# LEFT COLUMN — 3 sekce: Motivace, Téma, Throughline tagline
# Motivace = problém (agenti píšou rychleji, než kontrolujeme JAK)
# Téma     = vysvětlení tématu BP (jak měřit, ladit instrukce)
# Tagline  = throughline jako zapamatovatelná pointa

# Motivace
add_textbox(s, LEFT_X, Inches(2.05), LEFT_W, Inches(0.35),
            [("Motivace", True)], font_size=16, color=FIS_GREEN)
add_textbox(s, LEFT_X, Inches(2.45), LEFT_W, Inches(1.10), [
    "Agenti píšou kód rychleji,",
    "než kontrolujeme, jak ho píšou.",
], font_size=16)

# Téma
add_textbox(s, LEFT_X, Inches(3.70), LEFT_W, Inches(0.35),
            [("Téma", True)], font_size=16, color=FIS_GREEN)
add_textbox(s, LEFT_X, Inches(4.10), LEFT_W, Inches(1.20), [
    "Jak měřit chování agenta",
    "a podle dat ladit jeho instrukce.",
], font_size=16)

# Throughline tagline — krátký slogan, FIS green text na light gray pozadí
add_textbox(s, LEFT_X, Inches(5.60), LEFT_W, Inches(0.90), [
    [("Měříme agenta, ladíme instrukce.", True)],
], font_size=20, color=FIS_GREEN, fill=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# RIGHT COLUMN — Struktura prezentace (6 položek, žádné em dashe)
add_textbox(s, RIGHT_X, Inches(2.05), RIGHT_W, Inches(0.35),
            [("Struktura prezentace", True)], font_size=16, color=FIS_GREEN)

structure_items = [
    "1.  Cíle a metodika",
    "2.  Stav poznání",
    "3.  Sada metrik (cíl 1)",
    "4.  Iterativní ladění (cíl 2)",
    "5.  Ablace složek (cíl 3)",
    "6.  Závěr a přínos",
]
for i, item in enumerate(structure_items):
    add_textbox(s, RIGHT_X, Inches(2.55 + i * 0.50), RIGHT_W, Inches(0.45),
                [item], font_size=14)

set_notes(s, """[~75 s, čtyři sekce: Motivace, Téma, Throughline, Struktura prezentace]

[1] MOTIVACE (~20 s). AI agenti dnes umí napsat velké kusy kódu za desítky minut. Když takhle delegujeme víc práce, posouvá se klíčová otázka. Otázka už není, jestli kód projde testy. Otázka je, jestli kontrolujeme, JAK ho píše. Procesní stopa v repu (issue, větve, commity, pull requesty) konzistentně chybí, a benchmarky to neměří, protože měří jen funkčnost.

[2] TÉMA (~20 s). Tato práce se ptá, jak chování agenta měřit a podle těch dat ladit jeho instrukce. Tedy sada metrik, která pokrývá proces, produkt i zdroje, a postup, jak ji použít k iterativnímu vylepšování toho, jak agent pracuje.

[3] THROUGHLINE #1 (~10 s). Stručně: MĚŘÍME AGENTA, LADÍME INSTRUKCE. Tuto větu vyslovuji ještě dvakrát, na slidu 7 a 8.

[4] STRUKTURA PREZENTACE (~25 s). Půjdu tudy: cíle a metodika, krátký stav poznání. Hlavní část jsou dvě sady výsledků: iterativní ladění a ablace složek. Na závěr vlastní přínos a limity.""")

# -----------------------------------------------------------------------------
# SLIDE 3: Cíle a metodika
# -----------------------------------------------------------------------------
s = prs.slides[2]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, ["Cíle a metodika práce"])
    elif shp.placeholder_format and shp.placeholder_format.idx == 1:
        replace_text_keep_format(shp, [""])

remove_empty_placeholder(s, 1)

# Geometry matches title placeholder bounds for clean visual alignment
LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# === Cíle section ===
add_textbox(s, LEFT_X, Inches(2.05), CONTENT_W, Inches(0.35),
            [("Cíle", True)], font_size=16, color=FIS_GREEN)

add_textbox(s, LEFT_X, Inches(2.45), CONTENT_W, Inches(0.40),
            ["1.  Navrhnout sadu metrik pokrývající proces, produkt a zdroje."],
            font_size=14)
add_textbox(s, LEFT_X, Inches(2.90), CONTENT_W, Inches(0.40),
            ["2.  Ověřit proveditelnost iterativního ladění instrukcí podle těchto metrik."],
            font_size=14)
add_textbox(s, LEFT_X, Inches(3.35), CONTENT_W, Inches(0.40),
            ["3.  Ablacemi prozkoumat, které složky instrukcí přispívají a které jsou redundantní."],
            font_size=14)

# === Metodika section (Rámec / Postup / Měření — explicitní vrstvy) ===
add_textbox(s, LEFT_X, Inches(4.15), CONTENT_W, Inches(0.35),
            [("Metodika", True)], font_size=16, color=FIS_GREEN)

# Two-column layout for label/value alignment (proportional font would smear plain spaces)
LABEL_W = Inches(1.1)
VALUE_X = LEFT_X + LABEL_W
VALUE_W = CONTENT_W - LABEL_W
ROW_H = Inches(0.35)

# Row 1 — Rámec
add_textbox(s, LEFT_X, Inches(4.60), LABEL_W, ROW_H,
            [("Rámec:", True)], font_size=13)
add_textbox(s, VALUE_X, Inches(4.60), VALUE_W, ROW_H,
            ["Případová studie (Yin), systém upomínek faktur"], font_size=13)

# Row 2 — Postup (two visual rows: pilot + ablace)
add_textbox(s, LEFT_X, Inches(5.00), LABEL_W, ROW_H,
            [("Postup:", True)], font_size=13)
add_textbox(s, VALUE_X, Inches(5.00), VALUE_W, ROW_H,
            ["5 pilotních iterací r1–r5 (cíl 2)"], font_size=13)
add_textbox(s, VALUE_X, Inches(5.35), VALUE_W, ROW_H,
            ["+ 2 ablační varianty × 2 běhy (cíl 3)"], font_size=13)

# Row 3 — Měření
add_textbox(s, LEFT_X, Inches(5.75), LABEL_W, ROW_H,
            [("Měření:", True)], font_size=13)
add_textbox(s, VALUE_X, Inches(5.75), VALUE_W, ROW_H,
            ["sada metrik napříč procesem, produktem a zdroji"], font_size=13)

set_notes(s, """[~90 s]

Práce má tři cíle a opírá se o jednu případovou studii.

CÍL 1: navrhnout sadu metrik. Není to empirický běh, je to literárně a metodologicky odvozený návrh ze standardů kvality (Fenton a Bieman, ISO/IEC 25010). Detail rozeberu na slidu 4.

CÍL 2: ověřit, že iterativní ladění instrukcí funguje. Smyčka spustit, změřit, diagnostikovat, upravit, spustit znovu. Řeším pěti pilotními iteracemi r1 až r5. Mezi nimi se instrukce vyvíjejí podle dat, takže každá iterace je jiná konfigurace, ne replika.

CÍL 3: ablace. Z pilotu je vzata funkční konfigurace, jednotlivé složky se odeberou a měří se dopad. Dvě varianty, každá pouštěna dvakrát. Tady už jsou to repliky, kvůli stochasticitě modelu.

RÁMEC: Yin case study, jeden projekt do hloubky. Doména je systém upomínek faktur. Není to experiment ani DSR. Vědomě zůstávám na úrovni analytické generalizace, ne statistické. Cíl 2 proto formuluju jako PROVEDITELNOST, ne jako důkaz účinnosti.

MĚŘENÍ napříč třemi rovinami: proces, produkt, zdroje. Detail metrik a konkrétních výsledků na slidu 5.""")

# -----------------------------------------------------------------------------
# SLIDE 4: Stav poznání
# -----------------------------------------------------------------------------
s = prs.slides[3]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, ["Co se zatím měřilo a co chybí"])
    elif shp.placeholder_format and shp.placeholder_format.idx == 1:
        replace_text_keep_format(shp, [""])

remove_empty_placeholder(s, 1)

# Geometry aligned to title placeholder (same as slide 3)
LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# === Existující výzkum section ===
add_textbox(s, LEFT_X, Inches(2.05), CONTENT_W, Inches(0.35),
            [("Existující výzkum", True)], font_size=16, color=FIS_GREEN)

# Linie 1 — empirie (3 papers: METR 2026, Li 2026, Ehsani 2026)
add_textbox(s, LEFT_X, Inches(2.50), CONTENT_W, Inches(0.65),
            ["Tři studie shodně dokládají, že binární průchod testů nestačí: polovina prošlých změn by neprošla code review (METR 2026, Li 2026, Ehsani 2026)."],
            font_size=14)
# Linie 2 — benchmarky (klasické + novější Yin 2025)
add_textbox(s, LEFT_X, Inches(3.20), CONTENT_W, Inches(0.95),
            ["Funkční benchmarky (SWE-bench, Terminal-bench, HumanEval) měří binární průchod; novější vlna rozšiřuje na proces a tokeny (Yin et al. 2025), strukturální kvalita kódu zůstává mimo."],
            font_size=14)
# Linie 3 — studie souborů s instrukcemi
add_textbox(s, LEFT_X, Inches(4.20), CONTENT_W, Inches(0.65),
            ["Studie souborů s instrukcemi měří soubor jako celek (Lulla 2026: −28 % runtime; Gloaguen 2026: marginální přínos)."],
            font_size=14)

# === Mezera section ===
add_textbox(s, LEFT_X, Inches(5.05), CONTENT_W, Inches(0.35),
            [("Mezera", True)], font_size=16, color=FIS_GREEN)

add_textbox(s, LEFT_X, Inches(5.50), CONTENT_W, Inches(0.40),
            ["Holistický pohled napříč procesem, produktem a zdroji současně."],
            font_size=14)
add_textbox(s, LEFT_X, Inches(5.95), CONTENT_W, Inches(0.40),
            ["Izolace složek instrukcí, ne soubor jako blackbox."],
            font_size=14)

set_notes(s, """[~45 s]

Stav poznání má tři linie a jednu díru.

PRVNÍ LINIE: empirie. Tři nezávislé studie ukazují, že to, jestli kód projde testy, nestačí. METR 2026 analyzovali SWE-bench pull requesty a doložili, že polovina prošlých změn by neprošla code review. Li et al. potvrdili kvalitativní nedostatky na 1 210 mergedch pull requestech. Ehsani et al. na 33 tisících agentních pull requestech identifikovali procesní selhání jako významný podíl odmítnutí. Binární metrika tedy maskuje kvalitativní problém.

DRUHÁ LINIE: měřící rámce. Klasické benchmarky (SWE-bench, Terminal-bench, HumanEval) měří binárně. Novější vlna posunuje pohled dál: Yin et al. 2025 hodnotí 7 agentních frameworků na 3 osách (úspěch, efektivita reasoning, tokenová režie). To je posun, ale strukturální kvalita kódu (udržovatelnost, čistota) zůstává stranou.

TŘETÍ LINIE: studie souborů s instrukcemi. Lulla 2026 zkracuje běh o 28 procent, Gloaguen 2026 vidí jen marginální přínos. Obě studie ale měří soubor jako celek, neoddělují jeho složky.

MEZERA: Holistický pohled napříč procesem, produktem a zdroji současně, a izolace složek instrukcí. Sem tato práce cílí.""")

# -----------------------------------------------------------------------------
# Insert new slide AFTER slide 4 (Stav poznání) for Sada metrik (cíl 1)
# -----------------------------------------------------------------------------
insert_new_slide_after(prs, after_index=3, layout=content_layout)

# -----------------------------------------------------------------------------
# SLIDE 5: Sada metrik (cíl 1) — tři dimenze P/Q/E
# -----------------------------------------------------------------------------
s = prs.slides[4]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, ["Sada metrik: proces, produkt, zdroje"])
    elif shp.placeholder_format and shp.placeholder_format.idx == 1:
        replace_text_keep_format(shp, [""])
remove_empty_placeholder(s, 1)

# Geometry aligned to title placeholder
LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# Label nad obsahem (kotví slide na cíl 1)
add_textbox(s, LEFT_X, Inches(2.00), CONTENT_W, Inches(0.30),
            ["Cíl 1: sada metrik"],
            font_size=12, color=DARK_GRAY, italic=True)

# Intro věta
add_textbox(s, LEFT_X, Inches(2.40), CONTENT_W, Inches(0.40),
            ["Tři dimenze pokrývají, co binární průchod testů zakrývá."],
            font_size=14, color=DARK_GRAY)

# Three columns: PROCES / PRODUKT / ZDROJE
COL_W = Inches(2.736)
COL_GAP = Inches(0.15)
COL_P_X = LEFT_X                                # 0.688
COL_Q_X = LEFT_X + COL_W + COL_GAP              # 3.574
COL_E_X = LEFT_X + 2 * (COL_W + COL_GAP)        # 6.460

COL_TOP = Inches(3.00)

# Headers — barevně odlišené dimenze
add_textbox(s, COL_P_X, COL_TOP, COL_W, Inches(0.40),
            [("Proces (P1–P8)", True)],
            font_size=16, color=FIS_GREEN, align=PP_ALIGN.CENTER,
            fill=LIGHT_GRAY)
add_textbox(s, COL_Q_X, COL_TOP, COL_W, Inches(0.40),
            [("Produkt (Q1–Q8)", True)],
            font_size=16, color=FIS_GREEN, align=PP_ALIGN.CENTER,
            fill=LIGHT_GRAY)
add_textbox(s, COL_E_X, COL_TOP, COL_W, Inches(0.40),
            [("Zdroje (E1–E3)", True)],
            font_size=16, color=FIS_GREEN, align=PP_ALIGN.CENTER,
            fill=LIGHT_GRAY)

# Subtitle (italics, gray)
SUBTITLE_TOP = COL_TOP + Inches(0.45)
add_textbox(s, COL_P_X, SUBTITLE_TOP, COL_W, Inches(0.35),
            ["jak agent pracoval"],
            font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, COL_Q_X, SUBTITLE_TOP, COL_W, Inches(0.35),
            ["co vyrobil"],
            font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
add_textbox(s, COL_E_X, SUBTITLE_TOP, COL_W, Inches(0.35),
            ["kolik to stálo"],
            font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Examples per column (volný seznam, 4 příklady á 0.30" = 1.2")
EXAMPLES_TOP = SUBTITLE_TOP + Inches(0.45)
add_textbox(s, COL_P_X, EXAMPLES_TOP, COL_W, Inches(1.50), [
    "issues před kódem",
    "větev per issue",
    "test-first commity",
    "popisy PR a commitů",
], font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, COL_Q_X, EXAMPLES_TOP, COL_W, Inches(1.50), [
    "funkční korektnost",
    "kvalita testů",
    "lint, typy, složitost",
    "design (LLM soudce)",
], font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

add_textbox(s, COL_E_X, EXAMPLES_TOP, COL_W, Inches(1.50), [
    "tokeny",
    "čas běhu",
    "kompakce kontextu",
    "",
], font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Footer: rámec + měření (klidný, factual)
FOOTER_TOP = Inches(6.20)
add_textbox(s, LEFT_X, FOOTER_TOP, CONTENT_W, Inches(0.35),
            [[("Rámec: ", True), ("Fenton a Bieman (proces / produkt / zdroje), ISO/IEC 25010.", False)]],
            font_size=12, color=DARK_GRAY)
add_textbox(s, LEFT_X, FOOTER_TOP + Inches(0.40), CONTENT_W, Inches(0.35),
            [[("Měření: ", True), ("deterministické nástroje + LLM-as-judge (GLM-5) u tří kvalitativních metrik.", False)]],
            font_size=12, color=DARK_GRAY)

set_notes(s, """[~75 s]

Cíl 1: navrhnout sadu metrik. Tady je výsledek. 19 metrik napříč třemi dimenzemi.

PROCES (P1 až P8): jak agent pracoval. Jestli zakládal issues před kódem, jestli každý issue dostal svou větev, jestli psal testy před implementací, jestli psal popisy commitů a pull requestů. Tedy věci, které deterministické benchmarky nevidí, ale které rozhodují, jestli na ten kód v týmu navážeš.

PRODUKT (Q1 až Q8): co agent vyrobil. Funkční korektnost (projde proti referenční sadě?), kvalita testů (pokrytí, mutační skóre), deterministická kvalita kódu (lint, typecheck, složitost), a designová kvalita posuzovaná LLM soudcem.

ZDROJE (E1 až E3): kolik to stálo. Tokeny, čas, kompakce kontextu.

Rámec čerpá z Fenton a Bieman (rozdělení proces/produkt/zdroje) a ISO/IEC 25010 pro charakteristiky kvality kódu. U tří kvalitativních metrik, kde deterministický nástroj nestačí, je použit LLM-as-judge (GLM-5, jiná modelová rodina než hodnocený agent, kvůli self-enhancement bias).

Detail celé sady metrik je v záložním slidu na konci pro Q&A.""")

# -----------------------------------------------------------------------------
# Insert new slide AFTER slide 6 (Výsledky I) for Výsledky II (ablace)
# -----------------------------------------------------------------------------
insert_new_slide_after(prs, after_index=5, layout=content_layout)

# -----------------------------------------------------------------------------
# SLIDE 6: Výsledky I — iterace (cíl 2, operacionalizační oblouk z thesis)
# -----------------------------------------------------------------------------
s = prs.slides[5]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, [
            "Vzorec: pravidlo → příkaz → verifikační krok"
        ])
    elif shp.placeholder_format and shp.placeholder_format.idx == 1:
        replace_text_keep_format(shp, [""])

remove_empty_placeholder(s, 1)

# Geometry aligned to title placeholder
LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# Label nad obsahem (kotví slide na cíl 2)
add_textbox(s, LEFT_X, Inches(2.00), CONTENT_W, Inches(0.30),
            ["Cíl 2: iterativní postup"],
            font_size=12, color=DARK_GRAY, italic=True)

# === TikZ figure: Operacionalizační oblouk (z kap05.tex fig:operacionalizace) ===
# Slide nese jen jeden finding — operacionalizační oblouk. Konkrétní čísla
# (r3 hrdina, r4/r5 nuance) jdou do mluveného textu, ne na slide.
# Aspect ratio PNG ≈ 1.83:1 (1452 × 800), tj. při width=7.5" je height ≈ 4.1"
IMG_PATH = OUT.parent / "operacionalizace-1.png"
TIKZ_W = Inches(7.5)
TIKZ_LEFT = Inches((10 - 7.5) / 2)   # vodorovně vystředit (1.25")
s.shapes.add_picture(str(IMG_PATH), TIKZ_LEFT, Inches(2.70), width=TIKZ_W)

set_notes(s, """[~2:00, CENTERPIECE]

[SETUP ~20 s] Pilotní fáze: pět běhů, mezi nimi ladím instrukce podle metrik. Po dvou cyklech vidím opakující se vzorec.

[OBLOUK ~50 s] Tři nezávislá selhání (P2 větev per issue, P3 test-first commit, Q8 dokumentace API) prošla STEJNOU TRAJEKTORIÍ.

V r1 jsou v instrukcích jako PRAVIDLA: „každý issue má svou větev". Agent je ignoruje.

V r2 přepíšu na konkrétní PŘÍKAZY: „spusť git checkout -b issue-N". Agent je dodržuje formálně, ne podle záměru. Vytvoří jednu větev a do ní napíše tři issues.

V r3 přepíšu na VERIFIKAČNÍ KROK: „spusť `gh issue list --state open` a zkontroluj". A teprve tady agent dodržuje záměr.

[r3 HRDINA ~25 s] R3 je výsledkem tří takových úprav. Procesní checklist P1–P5 splněn celý. Q2 41 ze 42 referenčních testů. Q8 plně. A (což mě překvapilo) byl to NEJLEVNĚJŠÍ BĚH, asi sedmnáct centů. Pro porovnání: nejdražší pilot r4 stál šedesát čtyři centů a měl horší výsledky. Vyšší spotřeba tokenů nekorelovala s lepším výsledkem.

[NUANCE r5 ~25 s] R5 ukazuje druhou stranu. Stejné instrukce plus jedna úprava. Agent procesní postup IGNOROVAL ÚPLNĚ. Žádné issues, jeden monolitický commit. Ale jednokrokové verifikační příkazy (lint, typecheck) drží i tady. Q5 a Q7 nuly. Verifikační kroky nesou váhu i tehdy, když procesní sekvence selže.""")

# -----------------------------------------------------------------------------
# SLIDE 7: Výsledky II — ablace (cíl 3, HEAVY)
# -----------------------------------------------------------------------------
s = prs.slides[6]  # the inserted slide (druhý insert)
# The new slide from layout has placeholders: fill title, remove content placeholder
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, [
            "Verifikace drží kvalitu. Konvence drží design."
        ])
remove_empty_placeholder(s, 1)

LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# Label nad obsahem (kotví slide na cíl 3)
add_textbox(s, LEFT_X, Inches(2.00), CONTENT_W, Inches(0.30),
            ["Cíl 3: ablace složek"],
            font_size=12, color=DARK_GRAY, italic=True)

# === Horizontální flow: per ablace jeden řádek zleva doprava ===
# [Ablace X: co odebráno]  →  [Výsledek + interpretace]
# Konkrétní hodnoty (Q2 41→35, B-1 vs B-2 = 37/11) zůstávají v notes.

HEADER_W = Inches(3.30)
ARROW_W = Inches(0.40)
CLAIM_W = Inches(4.50)
GAP = Inches(0.15)
HEADER_X = LEFT_X                                          # 0.688
ARROW_X = HEADER_X + HEADER_W + GAP                        # 4.138
CLAIM_X = ARROW_X + ARROW_W + GAP                          # 4.688

# === Řádek 1: Ablace A — verifikace (předvídatelný) ===
# Šipky → vrácené: vizuálně oddělují „co odebráno" od „co se stalo".
# Pravá strana = claim (FIS green) + odůvodnění z thesis (italic gray).
ROW_A_TOP = Inches(2.50)
ROW_A_H = Inches(1.30)
add_textbox(s, HEADER_X, ROW_A_TOP, HEADER_W, ROW_A_H, [
    [("Ablace A", True, ACCENT_RED)],
    [("verifikační kroky odebrány.", False, DARK_GRAY)],
], font_size=14)

# Šipka — zarovnaná svisle na úroveň claimu
add_textbox(s, ARROW_X, ROW_A_TOP, ARROW_W, Inches(0.50),
            ["→"], font_size=22, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Claim
add_textbox(s, CLAIM_X, ROW_A_TOP, CLAIM_W, Inches(0.45),
            [[("Kvalita kódu klesla.", True, FIS_GREEN)]],
            font_size=15)

# Odůvodnění (z kap05.tex sec:vyhodnoceni-cil3)
add_textbox(s, CLAIM_X, ROW_A_TOP + Inches(0.55), CLAIM_W, Inches(0.80), [
    [("Předvídatelné: ", True),
     ("tyto verifikace byly do instrukcí přidány v pilotech právě proto, že bez nich metriky selhávaly.", False)],
], font_size=12, color=DARK_GRAY, italic=True)

# === Separator mezi A a B (tenká světle šedá linka) ===
SEP_Y = Inches(3.95)
sep = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                             LEFT_X, SEP_Y,
                             LEFT_X + CONTENT_W, SEP_Y)
sep.line.color.rgb = LIGHT_GRAY
sep.line.width = Pt(0.75)

# === Řádek 2: Ablace B — konvence (překvapivý) ===
ROW_B_TOP = Inches(4.20)
ROW_B_H = Inches(1.90)
add_textbox(s, HEADER_X, ROW_B_TOP, HEADER_W, ROW_B_H, [
    [("Ablace B", True, ACCENT_ORANGE)],
    [("kódové konvence odebrány.", False, DARK_GRAY)],
], font_size=14)

# Šipka — zarovnaná svisle, centrovaná vertikálně k 2-řádkovému claimu
add_textbox(s, ARROW_X, ROW_B_TOP + Inches(0.20), ARROW_W, Inches(0.50),
            ["→"], font_size=22, color=DARK_GRAY, align=PP_ALIGN.CENTER)

# Claim (2 řádky)
add_textbox(s, CLAIM_X, ROW_B_TOP, CLAIM_W, Inches(0.85), [
    [("Deterministika beze změny.", True, FIS_GREEN)],
    [("Klesla jen designová kvalita.", True, FIS_GREEN)],
], font_size=15)

# Odůvodnění
add_textbox(s, CLAIM_X, ROW_B_TOP + Inches(0.95), CLAIM_W, Inches(0.90), [
    [("Překvapivé: ", True),
     ("konvence drží z tréninku modelu, ne z instrukcí. Změnu zachytil jen LLM judge. Deterministické nástroje ji nevidí.", False)],
], font_size=12, color=DARK_GRAY, italic=True)

# === Poznámka pod čarou (nahrazuje šedý rámeček nuance) ===
add_textbox(s, LEFT_X, Inches(6.50), CONTENT_W, Inches(0.40),
            ["* Závěry jsou indikativní: variabilita mezi běhy byla srovnatelná s efektem ablace samotné."],
            font_size=11, color=DARK_GRAY, italic=True)

set_notes(s, """[~90 s]

[SETUP ~15 s] Z r3, vítězného běhu, byly provedeny dvě ablace, každá ve dvou bězích. Cílem bylo rozlišit, které složky AGENTS.md skutečně přispívaly.

[ABLACE A ~30 s] V ablaci A jsou odebrány VERIFIKAČNÍ KROKY: lint, typecheck, git log. Q2 klesla ze čtyřiceti jedné na třicet pět ze čtyřiceti dvou, lint warningy se ztrojnásobily. Bez explicitního příkazu agent eslint sám nepouští. Předvídatelný výsledek: tyhle verifikace byly do instrukcí přidány v pilotech právě proto, že bez nich metriky selhávaly.

[ABLACE B ~30 s] V ablaci B jsou odebrány KÓDOVÉ KONVENCE: modulární architektura, strict TypeScript, JSDoc. Deterministické metriky (lint, typy, složitost) zůstaly srovnatelné. Konvence drží z tréninku modelu, ne z instrukcí. Co spadlo, je Q8: designová kvalita posuzovaná LLM soudcem. To je překvapivé. Konvence ovlivňují rozhodnutí, která deterministické nástroje vůbec nezachytí.

[NUANCE ~15 s] Důležitá nuance: B-1 a B-2 měly Q2 sedmatřicet a jedenáct ze čtyřiceti dvou. Stejné instrukce, dramaticky rozdílný výsledek. Variabilita mezi běhy je srovnatelná s efektem ablace samotné. Dva běhy proto neumožňují kauzální závěr. ABLACE JSOU INDIKATIVNÍ.""")

# -----------------------------------------------------------------------------
# SLIDE 8: Závěr — sloučený přínos + limity
# -----------------------------------------------------------------------------
# Tři páry claim ↔ limit ve dvou sloupcích — paralelní struktura = honest hedge.
# Throughline „Měříme agenta, ladíme instrukce." na slidu NENÍ — řečník ji říká
# nahlas v mluveném textu, slide ji nepotřebuje opakovat.
s = prs.slides[7]
for shp in s.shapes:
    if not shp.has_text_frame:
        continue
    if shp.placeholder_format and shp.placeholder_format.idx == 0:
        replace_text_keep_format(shp, ["Závěr"])
remove_empty_placeholder(s, 1)

LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)
COL_GAP = Inches(0.20)
COL_W = Inches(4.154)
COL_L_X = LEFT_X
COL_R_X = LEFT_X + COL_W + COL_GAP

# === Headers ===
HEADERS_TOP = Inches(2.30)
add_textbox(s, COL_L_X, HEADERS_TOP, COL_W, Inches(0.40),
            [("Přínos", True)],
            font_size=18, color=FIS_GREEN)
add_textbox(s, COL_R_X, HEADERS_TOP, COL_W, Inches(0.40),
            [("Limity", True)],
            font_size=18, color=ACCENT_RED)

# === Tři páry claim ↔ limit ===
BODY_TOP = Inches(2.95)
ROW_H = Inches(0.95)   # větší výška pro dýchání mezi řádky

# Pár 1 — metrika ↔ scope
add_textbox(s, COL_L_X, BODY_TOP, COL_W, ROW_H,
            ["Sada metrik napříč procesem, produktem a zdroji."],
            font_size=15, color=DARK_GRAY)
add_textbox(s, COL_R_X, BODY_TOP, COL_W, ROW_H,
            ["Případová studie: jeden model, jeden projekt."],
            font_size=15, color=DARK_GRAY)

# Pár 2 — postup ↔ statistická síla
add_textbox(s, COL_L_X, BODY_TOP + ROW_H, COL_W, ROW_H,
            ["Iterativní postup po operacionalizačním oblouku."],
            font_size=15, color=DARK_GRAY)
add_textbox(s, COL_R_X, BODY_TOP + ROW_H, COL_W, ROW_H,
            ["Indikativní zjištění, ne statistický důkaz."],
            font_size=15, color=DARK_GRAY)

# Pár 3 — punchline ↔ nedeterminismus
add_textbox(s, COL_L_X, BODY_TOP + 2 * ROW_H, COL_W, ROW_H,
            [[("Přenositelný je postup, ne AGENTS.md.", True)]],
            font_size=15, color=DARK_GRAY)
add_textbox(s, COL_R_X, BODY_TOP + 2 * ROW_H, COL_W, ROW_H,
            ["Nedeterminismus modelu silný."],
            font_size=15, color=DARK_GRAY)

# === Poznámka pod čarou ===
add_textbox(s, LEFT_X, Inches(6.70), CONTENT_W, Inches(0.40),
            ["* Demonstrace proveditelnosti na MiniMax-M2.5 a TypeScript projektu. Aplikace na jiné modely a jazyky zůstává otevřená."],
            font_size=10, color=DARK_GRAY, italic=True)

set_notes(s, """[~60 s — CLOSER]

[PŘÍNOS ~30 s] Co práce přináší. Za prvé: sada metrik napříč procesem, produktem, zdroji. Rozlišuje typy selhání, které jediné výsledkové skóre zakryje. Za druhé: iterativní postup po operacionalizačním oblouku — pravidlo, příkaz, verifikační krok — replikoval se ve třech nezávislých nálezech. Za třetí (a klíčový bod): přenositelný je tenhle postup a sada metrik, ne konkrétní AGENTS.md. Pravidla vyladěná pro MiniMax-M2.5 na TypeScript projektu pro jiný model a jiný jazyk nemusí platit.

[LIMITY ~20 s] Kde to končí. Případová studie: jeden model, jeden projekt, malé n. Závěry jsou indikativní, ne kauzální: dva běhy na variantu nedovolí oddělit efekt ablace od nedeterminismu modelu. A nedeterminismus byl pozorován silný: v ablaci B-1 a B-2 dávaly stejné instrukce 37 a 11 ze 42 referenčních testů. Dramaticky rozdílný výsledek.

[CLOSER ~10 s] Měříme agenta, ladíme instrukce. Instrukce pro AI coding agenty se dají navrhovat jako měřená proměnná, ne psát od oka.

Děkuji vám za pozornost.

[Slide drží během Q&A.]""")

# -----------------------------------------------------------------------------
# Slide 9 (původní Závěr) sloučen do slidu 8 — odstranit placeholder.
# Tím se původní slide 10 (Otázky) posune na index 8.
# -----------------------------------------------------------------------------
remove_slide(prs, 8)

# -----------------------------------------------------------------------------
# SLIDE 9: Otázky — dle FIS šablony s otázkami z posudků
# -----------------------------------------------------------------------------
s = prs.slides[8]
remove_empty_placeholder(s, 1)

LEFT_X = Inches(0.688)
CONTENT_W = Inches(8.508)

# === Title ===
for shp in s.shapes:
    if shp.is_placeholder and shp.placeholder_format.idx == 0 and shp.has_text_frame:
        replace_text_keep_format(shp, ["Otázky"])
        break
else:
    add_textbox(s, LEFT_X, Inches(1.20), CONTENT_W, Inches(0.60),
                [("Otázky", True)],
                font_size=36, color=FIS_GREEN, align=PP_ALIGN.CENTER)

# === Dva sloupce: vedoucí vlevo, oponent vpravo ===
COL_W = Inches(3.95)
COL_L_X = LEFT_X
COL_R_X = Inches(5.20)
HDR_Y = Inches(2.00)
HDR_H = Inches(0.35)
Q_START_Y = Inches(2.60)
Q_H = Inches(0.65)

# Levý sloupec — vedoucí (4 otázky)
add_textbox(s, COL_L_X, HDR_Y, COL_W, HDR_H,
            [("Otázky vedoucího", True)],
            font_size=16, color=FIS_GREEN)
add_textbox(s, COL_L_X, Inches(2.35), COL_W, Inches(0.25),
            ["Ing. Jiří Korčák"],
            font_size=11, color=DARK_GRAY, italic=True)

for i, q in enumerate([
    "1. Role ablační studie — proč nestačily pilotní iterace?",
    "2. Rozdíl procesních, produktových a zdrojových metrik — konkrétní příklad?",
    "3. Jak odlišit vlastní přínos od přínosu AI asistenta?",
    "4. Pravidlo → příkaz → verifikační krok — konkrétní příklad?",
]):
    add_textbox(s, COL_L_X, Q_START_Y + i * Q_H, COL_W, Q_H,
                [q], font_size=13, color=DARK_GRAY)

# Pravý sloupec — oponent (1 otázka)
add_textbox(s, COL_R_X, HDR_Y, COL_W, HDR_H,
            [("Otázky oponenta", True)],
            font_size=16, color=FIS_GREEN)
add_textbox(s, COL_R_X, Inches(2.35), COL_W, Inches(0.25),
            ["Ing. Richard Antonín Novák, Ph.D."],
            font_size=11, color=DARK_GRAY, italic=True)

add_textbox(s, COL_R_X, Q_START_Y, COL_W, Q_H,
            ["1. Termín \"Ablace\" — hlubší vysvětlení pojmu a jeho role v práci?"],
            font_size=13, color=DARK_GRAY)

# Throughline na tomto slidu nepatří (ponechán jen na slidu 2 a 8).


set_notes(s, """[Q&A slide — drží během diskuse.]

=== OTÁZKY Z POSUDKŮ ===

VEDOUCÍ (Ing. Jiří Korčák) — 4 otázky:

1) Role ablační studie — proč nestačily piloty?
   → Ablace izolují jednotlivé složky instrukcí. Piloty mění celou konfiguraci najednou, takže nelze říct, která složka přispěla. Ablace = kontrola přičinění.

2) Rozdíl P/Q/E metrik + konkrétní příklad?
   → P = proces (issue tracking, commit messages, branch naming), Q = kvalita kódu (lint, typy, komplexita, LLM judge), E = zdroje (tokeny, čas, iterace).
   → Příklad: P1 (issue coverage) měří, jestli agent vytváří issues pro tasky. Q2 (lint warnings) měří čistotu kódu. E1 (token cost) měří cenu běhu.

3) Jak rozlišit tvůj přínos vs AI asistent?
   → AI asistent pomáhal s textem a kódem, ale design metrik, interpretace výsledků a metodologická rozhodnutí (proč case study, ne experiment) byla autorská. Transparentně deklarováno.

4) Pravidlo → příkaz → verifikační krok — konkrétní příklad?
   → Pravidlo: "kód má být kvalitní". Příkaz: "používej strict TypeScript". Verifikační krok: "spusť tsc --noEmit před commitem".
   → Bez verifikačního kroku agent příkaz ignoruje (nebo dělá jen někdy). S ním má explicitní checkpoint.

OPONENT (Ing. Richard Antonín Novák, Ph.D.) — 1 otázka:

1) Termín "Ablace" — hlubší vysvětlení?
   → Ablace = systematické odebírání složek z funkční konfigurace a měření dopadu. Analogie z ML: odebereš vrstvu z modelu a měříš, co se rozbije.
   → V mé práci: vezmu vítězný běh r3, odeberu verifikace (A) nebo konvence (B), spustím 2× a měřím metriky. Izolace přičinění.

=== ANTICIPACE DALŠÍCH OTÁZEK ===

- Proč tyto metriky? → Vyřazeno: příliš subjektivní (readability score), příliš drahé (mutation testing), neaplikovatelné (security audit).
- Proč MiniMax-M2.5? → Open-weight, dostupný, reprodukovatelný. GPT-4 by neumožnil ablace (API nedeterminismus + cena).
- Nejslabší místo? → Nedeterminismus: B-1 vs B-2 měly 37 vs 11 z 42 testů. Dva běhy nestačí na kauzální závěr.
- Co by udělal kolega? → Vzal by sadu metrik + oblouk pravidlo→příkaz→verifikace a aplikoval na svůj projekt. AGENTS.md by napsal sám podle svého modelu.""")

# -----------------------------------------------------------------------------
# SAVE
# -----------------------------------------------------------------------------
prs.save(str(OUT))
print(f"OK: written {OUT}")
print(f"Slides: {len(prs.slides)}")
for i, sl in enumerate(prs.slides, 1):
    title = ""
    for shp in sl.shapes:
        if shp.has_text_frame and shp.is_placeholder and shp.placeholder_format.idx == 0:
            title = shp.text_frame.text[:60]
            break
    print(f"  {i:2d}. {title}")
